"""Document processing system for various file formats."""

import hashlib
import mimetypes
from typing import Dict, List, Any, Optional, BinaryIO
from pathlib import Path
import io

from loguru import logger

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

from app.config import settings
from app.models.database import Document, DocumentChunk
from app.core.memory import memory


class DocumentProcessor:
    """Process various document types and extract text content."""
    
    def __init__(self):
        self.supported_formats = {
            '.pdf': self._process_pdf,
            '.docx': self._process_docx,
            '.doc': self._process_docx,  # Limited support
            '.pptx': self._process_pptx,
            '.xlsx': self._process_xlsx,
            '.xls': self._process_xlsx,  # Limited support
            '.txt': self._process_text,
            '.md': self._process_text,
            '.py': self._process_text,
            '.js': self._process_text,
            '.html': self._process_text,
            '.css': self._process_text,
            '.json': self._process_text,
            '.xml': self._process_text,
            '.csv': self._process_csv,
        }
    
    async def process_document(
        self, 
        file_path: str, 
        file_content: bytes, 
        original_filename: str,
        conversation_id: Optional[int] = None
    ) -> Dict[str, Any]:
        """Process a document and extract content."""
        try:
            # Determine file type
            file_extension = Path(original_filename).suffix.lower()
            content_type = mimetypes.guess_type(original_filename)[0]
            
            # Generate content hash
            content_hash = hashlib.sha256(file_content).hexdigest()
            
            # Check if document already exists
            existing_doc = self._check_existing_document(content_hash)
            if existing_doc:
                logger.info(f"Document already processed: {content_hash}")
                return existing_doc
            
            # Extract text content
            if file_extension in self.supported_formats:
                processor = self.supported_formats[file_extension]
                extracted_content = await processor(file_content, original_filename)
            else:
                logger.warning(f"Unsupported file format: {file_extension}")
                extracted_content = {
                    'text': f"Unsupported file format: {file_extension}",
                    'metadata': {'error': 'Unsupported format'}
                }
            
            # Create document record
            document_data = {
                'filename': Path(file_path).name,
                'original_filename': original_filename,
                'file_path': file_path,
                'file_size': len(file_content),
                'file_type': content_type or 'application/octet-stream',
                'content_hash': content_hash,
                'conversation_id': conversation_id,
                'processed': True,
                'metadata': {
                    'extension': file_extension,
                    'processing_method': processor.__name__ if file_extension in self.supported_formats else 'none',
                    **extracted_content.get('metadata', {})
                }
            }
            
            # Chunk the content for better processing
            chunks = self._chunk_content(
                extracted_content['text'], 
                chunk_size=1000, 
                overlap=100
            )
            
            # Store chunks in vector memory
            chunk_ids = []
            for i, chunk in enumerate(chunks):
                memory_id = memory.add_memory(
                    text=chunk['content'],
                    metadata={
                        'document_hash': content_hash,
                        'chunk_index': i,
                        'document_name': original_filename,
                        'document_type': file_extension,
                        'conversation_id': conversation_id
                    }
                )
                chunk_ids.append(memory_id)
            
            # Store document chunks info
            document_data['chunks'] = chunks
            document_data['chunk_ids'] = chunk_ids
            document_data['extracted_content'] = extracted_content
            
            logger.info(f"Successfully processed document: {original_filename}")
            return document_data
            
        except Exception as e:
            logger.error(f"Error processing document {original_filename}: {e}")
            return {
                'error': str(e),
                'filename': original_filename,
                'processed': False
            }
    
    async def _process_pdf(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Process PDF files."""
        if not PyPDF2:
            return {
                'text': 'PDF processing not available - PyPDF2 not installed',
                'metadata': {'error': 'Missing PyPDF2 dependency'}
            }
        
        try:
            pdf_file = io.BytesIO(file_content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            text_content = []
            metadata = {
                'page_count': len(pdf_reader.pages),
                'title': pdf_reader.metadata.get('/Title', '') if pdf_reader.metadata else '',
                'author': pdf_reader.metadata.get('/Author', '') if pdf_reader.metadata else '',
                'subject': pdf_reader.metadata.get('/Subject', '') if pdf_reader.metadata else '',
            }
            
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text.strip():
                        text_content.append(f"[Page {page_num + 1}]\n{page_text}")
                except Exception as e:
                    logger.warning(f"Error extracting text from page {page_num + 1}: {e}")
                    text_content.append(f"[Page {page_num + 1}]\n[Error extracting text]")
            
            return {
                'text': '\n\n'.join(text_content),
                'metadata': metadata
            }
            
        except Exception as e:
            logger.error(f"Error processing PDF {filename}: {e}")
            return {
                'text': f'Error processing PDF: {str(e)}',
                'metadata': {'error': str(e)}
            }
    
    async def _process_docx(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Process DOCX files."""
        if not DocxDocument:
            return {
                'text': 'DOCX processing not available - python-docx not installed',
                'metadata': {'error': 'Missing python-docx dependency'}
            }
        
        try:
            docx_file = io.BytesIO(file_content)
            doc = DocxDocument(docx_file)
            
            text_content = []
            metadata = {
                'paragraph_count': len(doc.paragraphs),
                'title': doc.core_properties.title or '',
                'author': doc.core_properties.author or '',
                'subject': doc.core_properties.subject or '',
                'created': str(doc.core_properties.created) if doc.core_properties.created else '',
                'modified': str(doc.core_properties.modified) if doc.core_properties.modified else '',
            }
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            # Extract table content
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text.strip())
                    table_text.append(' | '.join(row_text))
                if table_text:
                    text_content.append('\n'.join(table_text))
            
            return {
                'text': '\n\n'.join(text_content),
                'metadata': metadata
            }
            
        except Exception as e:
            logger.error(f"Error processing DOCX {filename}: {e}")
            return {
                'text': f'Error processing DOCX: {str(e)}',
                'metadata': {'error': str(e)}
            }
    
    async def _process_pptx(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Process PPTX files."""
        if not Presentation:
            return {
                'text': 'PPTX processing not available - python-pptx not installed',
                'metadata': {'error': 'Missing python-pptx dependency'}
            }
        
        try:
            pptx_file = io.BytesIO(file_content)
            prs = Presentation(pptx_file)
            
            text_content = []
            metadata = {
                'slide_count': len(prs.slides),
                'title': prs.core_properties.title or '',
                'author': prs.core_properties.author or '',
                'subject': prs.core_properties.subject or '',
            }
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = [f"[Slide {slide_num + 1}]"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if len(slide_text) > 1:  # More than just the slide header
                    text_content.append('\n'.join(slide_text))
            
            return {
                'text': '\n\n'.join(text_content),
                'metadata': metadata
            }
            
        except Exception as e:
            logger.error(f"Error processing PPTX {filename}: {e}")
            return {
                'text': f'Error processing PPTX: {str(e)}',
                'metadata': {'error': str(e)}
            }
    
    async def _process_xlsx(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Process XLSX files."""
        if not openpyxl:
            return {
                'text': 'XLSX processing not available - openpyxl not installed',
                'metadata': {'error': 'Missing openpyxl dependency'}
            }
        
        try:
            xlsx_file = io.BytesIO(file_content)
            workbook = openpyxl.load_workbook(xlsx_file, data_only=True)
            
            text_content = []
            metadata = {
                'sheet_count': len(workbook.worksheets),
                'sheet_names': [sheet.title for sheet in workbook.worksheets]
            }
            
            for sheet in workbook.worksheets:
                sheet_text = [f"[Sheet: {sheet.title}]"]
                
                # Get data from the sheet
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell in row:
                        if cell is not None:
                            row_text.append(str(cell))
                        else:
                            row_text.append('')
                    
                    if any(cell.strip() for cell in row_text):  # Skip empty rows
                        sheet_text.append(' | '.join(row_text))
                
                if len(sheet_text) > 1:  # More than just the sheet header
                    text_content.append('\n'.join(sheet_text))
            
            return {
                'text': '\n\n'.join(text_content),
                'metadata': metadata
            }
            
        except Exception as e:
            logger.error(f"Error processing XLSX {filename}: {e}")
            return {
                'text': f'Error processing XLSX: {str(e)}',
                'metadata': {'error': str(e)}
            }
    
    async def _process_text(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Process text-based files."""
        try:
            # Try different encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            text = None
            
            for encoding in encodings:
                try:
                    text = file_content.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            
            if text is None:
                # Fallback: decode with errors='replace'
                text = file_content.decode('utf-8', errors='replace')
            
            metadata = {
                'line_count': len(text.splitlines()),
                'character_count': len(text),
                'word_count': len(text.split()),
                'encoding': 'utf-8'  # We normalize to UTF-8
            }
            
            return {
                'text': text,
                'metadata': metadata
            }
            
        except Exception as e:
            logger.error(f"Error processing text file {filename}: {e}")
            return {
                'text': f'Error processing text file: {str(e)}',
                'metadata': {'error': str(e)}
            }
    
    async def _process_csv(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Process CSV files."""
        try:
            import csv
            
            # Decode content
            text = file_content.decode('utf-8', errors='replace')
            
            # Parse CSV
            csv_reader = csv.reader(io.StringIO(text))
            rows = list(csv_reader)
            
            # Format as text
            text_content = []
            for i, row in enumerate(rows):
                if i == 0:  # Header row
                    text_content.append("Headers: " + " | ".join(row))
                else:
                    text_content.append(" | ".join(row))
            
            metadata = {
                'row_count': len(rows),
                'column_count': len(rows[0]) if rows else 0,
                'has_header': True if rows else False
            }
            
            return {
                'text': '\n'.join(text_content),
                'metadata': metadata
            }
            
        except Exception as e:
            logger.error(f"Error processing CSV {filename}: {e}")
            return {
                'text': f'Error processing CSV: {str(e)}',
                'metadata': {'error': str(e)}
            }
    
    def _chunk_content(self, content: str, chunk_size: int = 1000, overlap: int = 100) -> List[Dict[str, Any]]:
        """Split content into overlapping chunks."""
        if not content:
            return []
        
        words = content.split()
        chunks = []
        
        for i in range(0, len(words), chunk_size - overlap):
            chunk_words = words[i:i + chunk_size]
            chunk_text = ' '.join(chunk_words)
            
            chunks.append({
                'index': len(chunks),
                'content': chunk_text,
                'word_count': len(chunk_words),
                'start_word': i,
                'end_word': min(i + chunk_size, len(words))
            })
        
        return chunks
    
    def _check_existing_document(self, content_hash: str) -> Optional[Dict[str, Any]]:
        """Check if document with same hash already exists."""
        # This would query the database for existing documents
        # For now, return None to always process
        return None
    
    async def search_documents(self, query: str, conversation_id: Optional[int] = None, k: int = 5) -> List[Dict[str, Any]]:
        """Search for relevant document content."""
        try:
            # Search in vector memory
            search_filters = {'document_type': True}  # Filter for document chunks
            if conversation_id:
                search_filters['conversation_id'] = conversation_id
            
            memories = memory.search_memories(query, k=k, threshold=0.6)
            
            # Filter for document memories
            document_memories = [
                mem for mem in memories 
                if mem.get('metadata', {}).get('document_name')
            ]
            
            return document_memories
            
        except Exception as e:
            logger.error(f"Error searching documents: {e}")
            return []
    
    def get_supported_formats(self) -> List[str]:
        """Get list of supported file formats."""
        return list(self.supported_formats.keys())


# Global document processor instance
document_processor = DocumentProcessor()